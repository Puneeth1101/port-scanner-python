# backend/document_processor/processor.py
import os
import tempfile
from typing import Dict, List, Optional, Tuple
import hashlib

# Document parsing libraries
import PyPDF2
import docx
import pptx
from pptx import Presentation
import csv
import json

class DocumentProcessor:
    """Handles parsing of different document types and chunking for embedding."""
    
    supported_extensions = {
        '.pdf': 'application/pdf',
        '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
        '.txt': 'text/plain',
        '.csv': 'text/csv',
        '.json': 'application/json',
    }
    
    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 200):
        """
        Initialize the document processor.
        
        Args:
            chunk_size: The size of text chunks for processing
            chunk_overlap: The overlap between chunks to maintain context
        """
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
    
    def process_document(self, file_path: str) -> List[Dict[str, str]]:
        """
        Process a document and return chunked text.
        
        Args:
            file_path: Path to the document
            
        Returns:
            List of dictionaries containing text chunks and metadata
        """
        _, ext = os.path.splitext(file_path)
        ext = ext.lower()
        
        if ext not in self.supported_extensions:
            raise ValueError(f"Unsupported file type: {ext}")
        
        # Extract text based on file type
        text = self._extract_text(file_path, ext)
        
        # Create document ID for tracking
        doc_id = self._generate_document_id(file_path)
        
        # Get file metadata
        metadata = self._extract_metadata(file_path)
        
        # Chunk the text
        chunks = self._chunk_text(text)
        
        # Create structured documents
        documents = []
        for i, chunk in enumerate(chunks):
            documents.append({
                'text': chunk,
                'metadata': {
                    'source': file_path,
                    'doc_id': doc_id,
                    'chunk_id': i,
                    'title': os.path.basename(file_path),
                    'file_type': ext[1:],  # Remove the dot
                    **metadata
                }
            })
        
        return documents
    def extract_text(file_input, filename=None):
        if isinstance(file_input, str):
        # it's a file path
            file_path = file_input
            ext = file_path.split('.')[-1]
        else:
        # it's a file-like object (BytesIO)
            ext = filename.split('.')[-1]

        if ext == 'pdf':
            import fitz
            doc = fitz.open(stream=file_input, filetype="pdf")
            return " ".join([page.get_text() for page in doc])
        elif ext == 'docx':
            from docx import Document
            if isinstance(file_input, str):
                doc = Document(file_input)
            else:
                doc = Document(file_input)
            return " ".join([para.text for para in doc.paragraphs])
        elif ext == 'pptx':
            from pptx import Presentation
            prs = Presentation(file_input)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return " ".join(text)
        return ""

    
    
    def _extract_from_pdf(self, file_path: str) -> str:
        """Extract text from PDF."""
        text = ""
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                text += page.extract_text() + "\n\n"
        return text
    
    def _extract_from_docx(self, file_path: str) -> str:
        """Extract text from DOCX."""
        doc = docx.Document(file_path)
        text = ""
        for para in doc.paragraphs:
            text += para.text + "\n"
        return text
    
    def _extract_from_pptx(self, file_path: str) -> str:
        """Extract text from PPTX."""
        presentation = Presentation(file_path)
        text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    
    def _extract_from_text(self, file_path: str) -> str:
        """Extract text from plain text file."""
        with open(file_path, 'r', encoding='utf-8', errors='replace') as file:
            return file.read()
    
    def _extract_from_csv(self, file_path: str) -> str:
        """Extract text from CSV."""
        text = ""
        with open(file_path, 'r', encoding='utf-8', errors='replace') as file:
            csv_reader = csv.reader(file)
            for row in csv_reader:
                text += " | ".join(row) + "\n"
        return text
    
    def _extract_from_json(self, file_path: str) -> str:
        """Extract text from JSON."""
        with open(file_path, 'r', encoding='utf-8', errors='replace') as file:
            data = json.load(file)
            return json.dumps(data, indent=2)
    
    def _chunk_text(self, text: str) -> List[str]:
        """Split text into chunks of specified size with overlap."""
        if not text:
            return []
            
        chunks = []
        start = 0
        text_length = len(text)
        
        while start < text_length:
            end = min(start + self.chunk_size, text_length)
            
            # If we're not at the end, try to find a good breaking point
            if end < text_length:
                # Look for paragraph break first
                paragraph_break = text.rfind('\n\n', start, end)
                if paragraph_break != -1 and paragraph_break > start + self.chunk_size // 2:
                    end = paragraph_break + 2  # Include the newlines
                else:
                    # Look for line break
                    line_break = text.rfind('\n', start, end)
                    if line_break != -1 and line_break > start + self.chunk_size // 2:
                        end = line_break + 1  # Include the newline
                    else:
                        # Look for sentence end (period followed by space)
                        sentence_end = text.rfind('. ', start, end)
                        if sentence_end != -1 and sentence_end > start + self.chunk_size // 2:
                            end = sentence_end + 2  # Include the period and space
                        else:
                            # Look for word boundary (space)
                            space = text.rfind(' ', start, end)
                            if space != -1 and space > start + self.chunk_size // 2:
                                end = space + 1  # Include the space
            
            chunks.append(text[start:end].strip())
            start = end - self.chunk_overlap
            
            # Prevent infinite loop on very short texts
            if start >= end:
                break
                
        return chunks
    
    def _generate_document_id(self, file_path: str) -> str:
        """Generate a unique document ID based on file path and modification time."""
        mod_time = os.path.getmtime(file_path)
        unique_str = f"{file_path}:{mod_time}"
        return hashlib.md5(unique_str.encode()).hexdigest()
    
    def _extract_metadata(self, file_path: str) -> Dict[str, str]:
        """Extract basic metadata from the file."""
        stats = os.stat(file_path)
        return {
            'size_bytes': stats.st_size,
            'created_at': stats.st_ctime,
            'modified_at': stats.st_mtime,
            'last_accessed': stats.st_atime,
        }